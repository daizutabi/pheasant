from pptx import Presentation
from pptx.util import Inches


def main():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    prs.save('test.pptx')

    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Adding a Table'

    rows = cols = 2
    left = top = Inches(2.0)
    width = Inches(6.0)
    height = Inches(0.8)

    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(2.0)
    table.columns[1].width = Inches(4.0)

    # write column headings
    table.cell(0, 0).text = 'Foo'
    table.cell(0, 1).text = 'Bar'

    # write body cells
    table.cell(1, 0).text = 'Baz'
    table.cell(1, 1).text = 'Qux'

    prs.save('test.pptx')
